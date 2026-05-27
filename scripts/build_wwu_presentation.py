"""
Generate the WWU presentation deck as a .pptx the team can touch up
before the meeting. Output: wwu_presentation_2026-05-27.pptx in the
project root.

DARK MODE EDITION: every surface, text color, table, and card flipped
to a dark slate palette so the deck is comfortable in any room and
under sunglasses. Brand red carries the accent role on every slide.

Visual palette:
  background      #0F1419 (deep slate near-black)
  surface         #1A2026 (card/table base)
  surface_high    #232A31 (alternating row)
  border          #2A323A (subtle divider)
  text primary    #F5F4F3 (off-white)
  text secondary  #B8B3AD (muted warm gray)
  text muted      #7C766F (deep muted)
  brand red       #F2483C (high-pop accent)
  red soft        #E07A6C (gentler red for prose)
  emerald         #10B981
  sky             #38BDF8
  violet          #A78BFA

Build: python scripts/build_wwu_presentation.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ---------- Dark palette ----------
BG = RGBColor(0x0F, 0x14, 0x19)
BG_TITLE = RGBColor(0x0A, 0x0F, 0x14)        # slightly darker for title slide
SURFACE = RGBColor(0x1A, 0x20, 0x26)
SURFACE_HIGH = RGBColor(0x23, 0x2A, 0x31)
BORDER = RGBColor(0x2A, 0x32, 0x3A)
TEXT_PRIMARY = RGBColor(0xF5, 0xF4, 0xF3)
TEXT_SECONDARY = RGBColor(0xB8, 0xB3, 0xAD)
TEXT_MUTED = RGBColor(0x7C, 0x76, 0x6F)
RED = RGBColor(0xF2, 0x48, 0x3C)
RED_SOFT = RGBColor(0xE0, 0x7A, 0x6C)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
SKY = RGBColor(0x38, 0xBD, 0xF8)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
AMBER = RGBColor(0xFB, 0xBF, 0x24)


# ---------- Deck setup ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def add_background(slide, color=BG):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill_solid(rect, color)
    no_line(rect)
    spTree = slide.shapes._spTree
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)
    return rect


def add_text(slide, text, *, left, top, width, height,
             size=18, bold=False, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT,
             font="Inter", italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_chip(slide, label, *, left, top, color, width=Inches(2.0),
             text_color=None):
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.44)
    )
    chip.adjustments[0] = 0.5
    fill_solid(chip, color)
    no_line(chip)
    tf = chip.text_frame
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Chivo Mono"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = text_color or TEXT_PRIMARY
    return chip


def add_accent_bar(slide, *, left, top, height, color=RED):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), height
    )
    fill_solid(bar, color)
    no_line(bar)
    return bar


def add_footer(slide, slide_number, total_slides):
    add_text(
        slide,
        "Civil Rights History Project    /    WWU    /    May 27, 2026",
        left=Inches(0.7), top=Inches(7.08),
        width=Inches(10), height=Inches(0.3),
        size=10, color=TEXT_MUTED, font="Chivo Mono",
    )
    add_text(
        slide,
        f"{slide_number:02d}  /  {total_slides:02d}",
        left=Inches(11.6), top=Inches(7.08),
        width=Inches(1.1), height=Inches(0.3),
        size=10, color=RED, font="Chivo Mono", align=PP_ALIGN.RIGHT, bold=True,
    )


def add_section_label(slide, label):
    add_text(
        slide, label.upper(),
        left=Inches(0.85), top=Inches(0.55),
        width=Inches(11), height=Inches(0.35),
        size=11, bold=True, color=RED, font="Chivo Mono",
    )


def add_title(slide, title):
    # Accent bar before the title
    add_accent_bar(
        slide,
        left=Inches(0.7), top=Inches(0.95), height=Inches(0.72),
        color=RED,
    )
    add_text(
        slide, title,
        left=Inches(0.95), top=Inches(0.9),
        width=Inches(12), height=Inches(1.0),
        size=32, bold=True, color=TEXT_PRIMARY, font="Inter",
    )


def add_bullets(slide, items, *, left, top, width, height,
                size=16, color=TEXT_PRIMARY, gap=0.6):
    """Render bullet items. Each item is (head, body) tuple."""
    cur_top = top
    for item in items:
        head, body = (item, "") if isinstance(item, str) else item
        # Red dot marker
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left, cur_top + Inches(0.13),
            Inches(0.16), Inches(0.16),
        )
        fill_solid(marker, RED)
        no_line(marker)

        box = slide.shapes.add_textbox(
            left + Inches(0.36), cur_top, width - Inches(0.36), Inches(gap + 0.2)
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run_head = p.add_run()
        run_head.text = head
        run_head.font.name = "Inter"
        run_head.font.size = Pt(size)
        run_head.font.bold = True
        run_head.font.color.rgb = color
        if body:
            run_body = p.add_run()
            run_body.text = "    " + body
            run_body.font.name = "Inter"
            run_body.font.size = Pt(size - 1)
            run_body.font.bold = False
            run_body.font.color.rgb = TEXT_SECONDARY
        cur_top += Inches(gap)


def add_table(slide, header, rows, *, left, top, width, height,
              first_col_color=RED):
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(header), left, top, width, height
    )
    table = table_shape.table
    for i, h in enumerate(header):
        cell = table.cell(0, i)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.name = "Chivo Mono"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = RED
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG
        cell.margin_left = Inches(0.12)
        cell.margin_right = Inches(0.12)
        cell.margin_top = Inches(0.06)
        cell.margin_bottom = Inches(0.06)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.name = "Inter"
            run.font.size = Pt(12)
            if c == 0:
                run.font.bold = True
                run.font.color.rgb = first_col_color
            else:
                run.font.color.rgb = TEXT_PRIMARY if r % 2 == 1 else TEXT_PRIMARY
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE if r % 2 == 1 else SURFACE_HIGH
            cell.margin_left = Inches(0.14)
            cell.margin_right = Inches(0.14)
            cell.margin_top = Inches(0.08)
            cell.margin_bottom = Inches(0.08)
    return table_shape


def add_card(slide, *, left, top, width, height,
             fill=SURFACE, border_color=BORDER, border_width=1):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_solid(card, fill)
    card.line.color.rgb = border_color
    card.line.width = Pt(border_width)
    return card


# ---------- Slides ----------
def slide_01_title():
    s = prs.slides.add_slide(BLANK)
    add_background(s, color=BG_TITLE)

    # Big decorative bar on the left
    add_accent_bar(s, left=Inches(0), top=Inches(0), height=prs.slide_height, color=RED)

    # Top section label
    add_text(
        s, "CIVIL RIGHTS HISTORY PROJECT",
        left=Inches(0.95), top=Inches(0.9),
        width=Inches(11), height=Inches(0.4),
        size=13, bold=True, color=RED, font="Chivo Mono",
    )

    # Main title
    add_text(
        s, "From Whisper",
        left=Inches(0.95), top=Inches(1.55),
        width=Inches(12), height=Inches(1.1),
        size=56, bold=True, color=TEXT_PRIMARY, font="Inter",
    )
    add_text(
        s, "to Smithsonian-Grade",
        left=Inches(0.95), top=Inches(2.45),
        width=Inches(12), height=Inches(1.1),
        size=56, bold=True, color=TEXT_PRIMARY, font="Inter",
    )
    add_text(
        s, "Publication.",
        left=Inches(0.95), top=Inches(3.35),
        width=Inches(12), height=Inches(1.1),
        size=56, bold=True, color=RED, font="Inter",
    )

    # Subtitle
    add_text(
        s,
        "Nine audit passes  /  Library of Congress as canonical authority  /  live RAG layer",
        left=Inches(0.95), top=Inches(4.6),
        width=Inches(12), height=Inches(0.5),
        size=17, color=TEXT_SECONDARY, font="Source Serif 4", italic=True,
    )

    # Divider
    line = s.shapes.add_connector(
        1, Inches(0.95), Inches(5.4), Inches(7.5), Inches(5.4)
    )
    line.line.color.rgb = RED
    line.line.width = Pt(2)

    # Footer info
    add_text(
        s, "WWU TEAM MEETING",
        left=Inches(0.95), top=Inches(5.7),
        width=Inches(11), height=Inches(0.4),
        size=12, bold=True, color=RED, font="Chivo Mono",
    )
    add_text(
        s, "May 27, 2026",
        left=Inches(0.95), top=Inches(6.05),
        width=Inches(11), height=Inches(0.4),
        size=14, color=TEXT_PRIMARY, font="Chivo Mono",
    )
    add_text(
        s, "civil-rights-staging.netlify.app",
        left=Inches(0.95), top=Inches(6.45),
        width=Inches(11), height=Inches(0.4),
        size=13, color=TEXT_MUTED, font="Chivo Mono",
    )
    return s


def slide_02_problem():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "The problem in one sentence")
    add_title(s, "Whisper transcripts read fluently.  They are also wrong.")
    add_text(
        s,
        "Automatic Speech Recognition on 1960s-era oral history audio produces "
        "transcript text that reads fluently but contains systematic errors invisible "
        "to a casual read. A Smithsonian-grade publication built on those transcripts "
        "would confidently say the wrong thing.",
        left=Inches(0.95), top=Inches(2.2),
        width=Inches(11.5), height=Inches(1.5),
        size=18, color=TEXT_SECONDARY, font="Source Serif 4",
    )

    # Two side-by-side cards
    card_w = Inches(5.65)
    card_h = Inches(2.7)
    card_y = Inches(4.1)
    # Left: Whisper output
    left_card = add_card(s, left=Inches(0.95), top=card_y, width=card_w, height=card_h)
    add_accent_bar(s, left=Inches(0.95), top=card_y, height=card_h, color=RED)
    add_text(
        s, "WHAT WHISPER SAID",
        left=Inches(1.25), top=card_y + Inches(0.2),
        width=card_w - Inches(0.4), height=Inches(0.35),
        size=11, bold=True, color=RED, font="Chivo Mono",
    )
    add_text(
        s, '"Earl, Adam Clayton Powell Sr., Andrew, Carlos"',
        left=Inches(1.25), top=card_y + Inches(0.7),
        width=card_w - Inches(0.4), height=Inches(0.9),
        size=18, color=TEXT_PRIMARY, font="Source Serif 4", italic=True,
    )
    add_text(
        s, "John Carlos's siblings, in his own voice.",
        left=Inches(1.25), top=card_y + Inches(2.0),
        width=card_w - Inches(0.4), height=Inches(0.4),
        size=11, color=TEXT_MUTED, font="Chivo Mono",
    )

    # Right: Actual
    right_card = add_card(
        s, left=Inches(6.95), top=card_y, width=card_w, height=card_h
    )
    add_accent_bar(s, left=Inches(6.95), top=card_y, height=card_h, color=EMERALD)
    add_text(
        s, "WHAT JOHN CARLOS ACTUALLY SAID",
        left=Inches(7.25), top=card_y + Inches(0.2),
        width=card_w - Inches(0.4), height=Inches(0.35),
        size=11, bold=True, color=EMERALD, font="Chivo Mono",
    )
    add_text(
        s, '"Earl Jr., Andrew, John Carlos"',
        left=Inches(7.25), top=card_y + Inches(0.7),
        width=card_w - Inches(0.4), height=Inches(0.9),
        size=18, color=TEXT_PRIMARY, font="Source Serif 4", italic=True,
    )
    add_text(
        s,
        "Whisper bled Adam Clayton Powell Sr. (the Abyssinian "
        "Baptist Church pastor mentioned three paragraphs later) "
        "backward into the siblings list, inventing a fictional brother.",
        left=Inches(7.25), top=card_y + Inches(1.55),
        width=card_w - Inches(0.4), height=Inches(1.1),
        size=11, color=TEXT_SECONDARY, font="Source Serif 4",
    )
    return s


def slide_03_error_examples():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "Concrete examples from the corpus")
    add_title(s, "Categorical errors, not typos.")
    add_text(
        s,
        "Wrong person, wrong place, wrong fact, presented with the same confidence "
        "as the rest of the text.",
        left=Inches(0.95), top=Inches(2.0),
        width=Inches(11.5), height=Inches(0.6),
        size=15, color=TEXT_SECONDARY, font="Source Serif 4", italic=True,
    )
    rows = [
        ("David Klein", "David Cline", "Wrong interviewer attributed to 100+ interviews"),
        ('"Daniel H. Krenge De Iongh"', "Daniel H. Crena de Iongh", "First Treasurer of the World Bank; key apartheid-finance figure"),
        ('"Paul Hoffman Robeson"', "Paul Robeson + Paul Hoffman", "Two distinct figures merged into one; 4 occurrences"),
        ('"Margaret the King"', "Martin Luther King Jr.", "Recurring across multiple interviews"),
        ('"Auto bom barroom"', "Audubon Ballroom", "Site of Malcolm X's assassination; unrecognizable"),
        ('"Stoke and Carmichael"', "Stokely Carmichael", "SNCC chairman; Whisper invented a second speaker"),
        ('"Elders Cleaver"', "Eldridge Cleaver", "BPP Minister of Information; 24 occurrences in one transcript"),
    ]
    add_table(
        s,
        ["Whisper output", "Correct reading", "Why it matters"],
        rows,
        left=Inches(0.95), top=Inches(2.85),
        width=Inches(11.45), height=Inches(3.9),
    )
    return s


def slide_04_audit_cascade():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "The audit cascade")
    add_title(s, "Eight passes of cleaning  +  a ninth against the LoC archive.")
    passes = [
        ("Pass 1", "Phonetic alias matching + ground-truth corpus grounding."),
        ("Pass 2", "Per-entry tail sweep for entries with partial Pass 1 reads."),
        ("Pass 3", "Consolidation: confidence resolutions + adversarial-review flags."),
        ("Pass 4", "Sweeping QA + fact-check; one transcript per agent, strict isolation."),
        ("Layer 5", "Corpus-global fidelity audit: phantom renderings, canonical consistency, contradictions."),
        ("Pass 6", "Apply-step discipline lockdown after the Pass 7 PRR side-file regression."),
        ("Pass 7 PRR", "Narrative-coherence pass: catches ASR name-bleed cases alias matching cannot."),
        ("Pass 8", "Library of Congress canonical-archive cross-reference. The first pass anchored to an external authority. Coverage: 127 / 127."),
    ]
    cur = Inches(2.1)
    for label, body in passes:
        # Choose chip color: Pass 8 in red for emphasis, the rest in a deep slate
        is_loc = label == "Pass 8"
        chip_color = RED if is_loc else SURFACE_HIGH
        text_color = TEXT_PRIMARY
        add_chip(
            s, label,
            left=Inches(0.95), top=cur,
            color=chip_color, width=Inches(1.5),
            text_color=text_color,
        )
        add_text(
            s, body,
            left=Inches(2.6), top=cur + Inches(0.07),
            width=Inches(10), height=Inches(0.5),
            size=14, color=TEXT_PRIMARY if is_loc else TEXT_SECONDARY, font="Inter",
            bold=is_loc,
        )
        cur += Inches(0.6)
    return s


def slide_05_loc_pass():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "Pass 8: Library of Congress healing")
    add_title(s, "The 9th pass that anchored the corpus to its canonical authority.")
    add_text(
        s,
        "LoC and the Smithsonian NMAAHC jointly produced this corpus. LoC has authoritative "
        "transcripts for every interview. Passes 1-7 used internal ground truth + adversarial "
        "review but never cross-referenced our text against LoC's. Pass 8 fixed that.",
        left=Inches(0.95), top=Inches(2.05),
        width=Inches(11.5), height=Inches(1.4),
        size=15, color=TEXT_SECONDARY, font="Source Serif 4", italic=True,
    )
    steps = [
        ("1.", "Resolve LoC item URL for each interviewee (search API + PDF fallback)."),
        ("2.", "Word-align our Whisper-derived .srt against LoC's transcript text."),
        ("3.", "Classify each divergence under conservative-first-pass discipline."),
        ("4.", "Apply only the deterministic-verdict heals inside existing cue boundaries."),
        ("5.", "Flag the rest for SME review (per-entry stage file: pass8_stage/...)."),
        ("6.", "Audit-canon safeguard: do not auto-reverse a confirmed prior-pass correction."),
    ]
    cur = Inches(3.55)
    for num, body in steps:
        add_text(
            s, num,
            left=Inches(0.95), top=cur,
            width=Inches(0.6), height=Inches(0.4),
            size=17, bold=True, color=RED, font="Chivo Mono",
        )
        add_text(
            s, body,
            left=Inches(1.55), top=cur + Inches(0.03),
            width=Inches(11), height=Inches(0.5),
            size=14, color=TEXT_PRIMARY, font="Inter",
        )
        cur += Inches(0.5)
    return s


def slide_06_coverage():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "Coverage achieved")
    add_title(s, "Every audit-able interview cross-validated against LoC.")

    # Hero stats row
    stats_y = Inches(2.1)
    stat_box_w = Inches(3.65)
    stat_h = Inches(1.4)

    for i, (big, small) in enumerate([
        ("100%", "of the 127 audit-able corpus healed"),
        ("127 / 127", "interviews cross-referenced with LoC"),
        ("~2,600", "ASR-error heals applied; 0 apply failures"),
    ]):
        x = Inches(0.95) + (stat_box_w + Inches(0.2)) * i
        add_card(s, left=x, top=stats_y, width=stat_box_w, height=stat_h)
        add_text(
            s, big,
            left=x + Inches(0.2), top=stats_y + Inches(0.1),
            width=stat_box_w - Inches(0.4), height=Inches(0.7),
            size=34, bold=True, color=RED, font="Inter",
        )
        add_text(
            s, small,
            left=x + Inches(0.2), top=stats_y + Inches(0.85),
            width=stat_box_w - Inches(0.4), height=Inches(0.5),
            size=11, color=TEXT_SECONDARY, font="Chivo Mono",
        )

    # Detail table
    rows = [
        ("Total interviews in audit-able corpus", "127"),
        ("Healed via LoC TEI2 XML transcripts", "92"),
        ("Healed via LoC PDF text extraction", "35"),
        ("Audio-only on LoC (no transcript at all)", "0"),
        ("Total divergences flagged for SME review", "~95,000"),
        ("Cue-count / timestamp verification failures", "0"),
    ]
    add_table(
        s,
        ["Metric", "Value"],
        rows,
        left=Inches(0.95), top=Inches(3.85),
        width=Inches(11.45), height=Inches(2.8),
    )
    add_text(
        s,
        "LoC has machine-extractable transcripts for every interview in our corpus.  Coverage is complete.",
        left=Inches(0.95), top=Inches(6.75),
        width=Inches(11.5), height=Inches(0.35),
        size=13, color=RED, font="Inter", bold=True, italic=True,
    )
    return s


def slide_07_conservative():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "Conservative-first-pass discipline")
    add_title(s, "95,000 divergences detected.   ~2,600 applied.   Why so few?")

    add_text(
        s,
        "The automatic-apply path is reserved for cases we are highly confident:",
        left=Inches(0.95), top=Inches(2.1),
        width=Inches(11.5), height=Inches(0.5),
        size=15, color=TEXT_SECONDARY, font="Source Serif 4", italic=True,
    )
    bullets = [
        ("Single-word proper-noun substitution.",
         "Klein -> Cline.  Phonetic confusion band only; anything broader skipped."),
        ("Similarity ratio in the 0.55 - 0.95 band.",
         "Close enough to be ASR confusion, far enough not to be case or orthography drift."),
        ("Our token NOT in the audit-canon set.",
         "Never reverse a confirmed prior-pass correction; surface as SME review instead."),
        ("Adjacent characters not contraction or hyphen.",
         "Avoids the landmines: 'don' -> 'Daniel H. Crena de Iongh' would eat every 'don't'."),
    ]
    add_bullets(
        s, bullets,
        left=Inches(0.95), top=Inches(2.85),
        width=Inches(11.5), height=Inches(3.0),
        size=15, gap=0.7,
    )

    # Footer box
    box = add_card(
        s, left=Inches(0.95), top=Inches(5.85),
        width=Inches(11.45), height=Inches(1.05),
        fill=SURFACE_HIGH, border_color=RED, border_width=2,
    )
    add_text(
        s,
        "The conservative path means we do not auto-introduce new errors while fixing old ones.  Everything outside the clean buckets is preserved verbatim and catalogued for human review.",
        left=Inches(1.15), top=Inches(5.98),
        width=Inches(11.05), height=Inches(0.85),
        size=13, color=TEXT_PRIMARY, font="Source Serif 4", italic=True,
    )
    return s


def slide_08_challenges():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "Challenge categories")
    add_title(s, "What else broke, and how we handled it.")
    rows = [
        ("Incorrect header / interviewer attribution",
         '"David Klein" everywhere it should have been David Cline.  Pass 8 fixed against LoC.'),
        ("Incomplete transcripts",
         "Whisper sometimes emptied joint interviews.  Re-ingested via the streamlined pipeline."),
        ("Multiple speakers in one cue",
         "ASR cannot segment overlapping voices.  Flagged as ingestion-only tier; retrieval hedges."),
        ("Corrupted transcriptions",
         "Audubon Ballroom -> Auto bom barroom.  Lenox -> Linux.  Pass 8 catches when LoC has the canonical form."),
        ("ASR name-bleed",
         "Two figures merged into one (Paul Robeson + Paul Hoffman).  Requires narrative-coherence pass."),
        ("Short-needle substitution corruption",
         "'don' -> 'Daniel H. Crena de Iongh' would eat 'don't'.  Word-boundary + contraction guards."),
        ("Audit-canon leakage",
         "Whisper spelling promoted in an early pass, never re-checked.  Pass 8 surfaces 710 such cases."),
    ]
    add_table(
        s,
        ["Challenge", "How it surfaces / how we handled it"],
        rows,
        left=Inches(0.95), top=Inches(2.05),
        width=Inches(11.45), height=Inches(4.6),
    )
    return s


def slide_09_what_it_enables():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "What the cleaned corpus enables")
    add_title(s, "A live, interactive RAG layer on top of an audit-grade substrate.")
    add_text(
        s,
        "Audit cleanliness is the foundation.  The retrieval layer is what makes the corpus "
        "useful to a researcher today.",
        left=Inches(0.95), top=Inches(2.0),
        width=Inches(11.5), height=Inches(0.7),
        size=15, color=TEXT_SECONDARY, font="Source Serif 4", italic=True,
    )
    bullets = [
        ("Voyage AI voyage-3 embeddings (1024-dim).",
         "Retrieval-tuned; finds passages that ANSWER a question, not just mention the same words."),
        ("Pinecone Builder serverless vector index.",
         "15,464 time-anchored .srt chunks across the 136-entry unified corpus."),
        ("Voyage rerank-2 on top of vector search.",
         "Reorders the initial 50-candidate pull by deeper semantic relevance."),
        ("5-tier fidelity badge on every passage.",
         "low / medium / publication-block / not-auditable / ingestion-only.  Smithsonian rigor at result level."),
        ("Library of Congress catalog deep-link on every result.",
         "The audience can verify the source against the canonical archive in one click."),
    ]
    add_bullets(
        s, bullets,
        left=Inches(0.95), top=Inches(2.9),
        width=Inches(11.5), height=Inches(4.0),
        size=14, gap=0.8,
    )
    return s


def slide_10_demos():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "/rag-explore   -   the interactive demo page")
    add_title(s, "Six surfaces on one page, all on the same embedding substrate.")
    demos = [
        ("Spectrum", VIOLET,
         "136 dots placed along one named conceptual axis at a time.  Nonviolence as Theology vs. Armed Self-Defense, Sacred vs. Secular Framing, etc.  Click a dot to drill in."),
        ("Semantic Overlap", SKY,
         "Pick an interview.  See which other voices in the corpus discuss semantically-related material.  The 'voices in conversation' demo."),
        ("Word Search", EMERALD,
         "Four 2D scatters across pairs of named axes plus a 5-spectrum 1D summary that lights up when you project a phrase onto every axis at once."),
        ("Interview Map", RED,
         "136 dots in UMAP-projected space with empirically-derived axis labels extracted from the corpus itself (Medical Law, Movement, Family, Crime)."),
        ("Quote Finder", AMBER,
         "Paste a half-remembered quote.  Get the source with the exact timestamp and Library of Congress catalog link."),
        ("Secondary surfaces", TEXT_MUTED,
         "Themes  /  Famous Names  /  Atlas  /  Network  /  Tours  /  Quote of the Day.  All precomputed JSON; no LLM per query."),
    ]
    grid_left = Inches(0.95)
    grid_top = Inches(2.05)
    card_w = Inches(5.65)
    card_h = Inches(1.55)
    gap_x = Inches(0.15)
    gap_y = Inches(0.15)
    for i, (label, color, body) in enumerate(demos):
        col = i % 2
        row = i // 2
        x = grid_left + (card_w + gap_x) * col
        y = grid_top + (card_h + gap_y) * row
        add_card(s, left=x, top=y, width=card_w, height=card_h)
        add_accent_bar(s, left=x, top=y, height=card_h, color=color)
        add_chip(
            s, label,
            left=x + Inches(0.25), top=y + Inches(0.2),
            color=color, width=Inches(2.4),
            text_color=BG if color == AMBER else TEXT_PRIMARY,
        )
        add_text(
            s, body,
            left=x + Inches(0.25), top=y + Inches(0.75),
            width=card_w - Inches(0.45), height=card_h - Inches(0.85),
            size=10.5, color=TEXT_SECONDARY, font="Source Serif 4",
        )
    return s


def slide_mcp():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "MCP server   -   the corpus as a connector")
    add_title(s, "Same substrate, different interface.  Reachable from any MCP client.")
    add_text(
        s,
        "mcp-server/server.mjs exposes the same Pinecone + Voyage substrate to LLM chat "
        "clients (Codex Desktop, Claude Desktop, the Anthropic Connector Directory).  "
        "Six tools, three of them research-pattern, all citation-grade.",
        left=Inches(0.95), top=Inches(2.0),
        width=Inches(11.5), height=Inches(1.3),
        size=14, color=TEXT_SECONDARY, font="Source Serif 4", italic=True,
    )

    # Two-column layout
    col_w = Inches(5.65)
    col_top = Inches(3.55)
    # Left column
    add_text(
        s, "PRIMITIVE TOOLS",
        left=Inches(0.95), top=col_top,
        width=col_w, height=Inches(0.35),
        size=11, bold=True, color=RED, font="Chivo Mono",
    )
    primitives = [
        ("search_transcripts",
         "Citation-grade semantic search.  Filters by entry, tier, dedupe-by-entry."),
        ("get_transcript",
         "Full ordered transcript for one entry_number."),
        ("list_leaders",
         "Archive roster: entry numbers, LoC URLs, provenance, audit tiers."),
    ]
    cur = col_top + Inches(0.4)
    for name, body in primitives:
        add_text(
            s, name,
            left=Inches(0.95), top=cur,
            width=col_w, height=Inches(0.32),
            size=13, bold=True, color=EMERALD, font="Chivo Mono",
        )
        add_text(
            s, body,
            left=Inches(0.95), top=cur + Inches(0.32),
            width=col_w, height=Inches(0.65),
            size=11, color=TEXT_SECONDARY, font="Source Serif 4",
        )
        cur += Inches(1.0)

    # Right column
    add_text(
        s, "RESEARCH-PATTERN TOOLS",
        left=Inches(6.75), top=col_top,
        width=col_w, height=Inches(0.35),
        size=11, bold=True, color=RED, font="Chivo Mono",
    )
    patterns = [
        ("compare_perspectives({ topic })",
         "Multiple interviewee voices on one topic, deduped, with citation framing."),
        ("trace_evolution({ interviewee, topic })",
         "Chronologically ordered passages from one interview, scoped to one topic."),
        ("source_for_claim({ claim })",
         "Passages that support, complicate, or contradict a claim.  Citation metadata preserved."),
    ]
    cur = col_top + Inches(0.4)
    for name, body in patterns:
        add_text(
            s, name,
            left=Inches(6.75), top=cur,
            width=col_w, height=Inches(0.32),
            size=13, bold=True, color=SKY, font="Chivo Mono",
        )
        add_text(
            s, body,
            left=Inches(6.75), top=cur + Inches(0.32),
            width=col_w, height=Inches(0.65),
            size=11, color=TEXT_SECONDARY, font="Source Serif 4",
        )
        cur += Inches(1.0)

    # Footer band
    band = add_card(
        s, left=Inches(0.95), top=Inches(6.55),
        width=Inches(11.45), height=Inches(0.5),
        fill=RED, border_color=RED, border_width=0,
    )
    add_text(
        s,
        "Every tool returns Chicago / APA / MLA citation blocks + the LoC catalog URL.   Deployable to Fly.io; verified locally against the live Pinecone index.",
        left=Inches(1.15), top=Inches(6.62),
        width=Inches(11.05), height=Inches(0.45),
        size=11, color=TEXT_PRIMARY, font="Chivo Mono", bold=True,
    )
    return s


def slide_11_philosophy():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "Philosophy of embedding")
    add_title(s, "When two voices land within 0.12 cosine, the embedding heard them.")
    add_text(
        s,
        "Cosine similarity is reading a thematic kinship the speakers may never have realized "
        "they shared.  This is the conceptual move that turns retrieval into discovery.",
        left=Inches(0.95), top=Inches(2.0),
        width=Inches(11.5), height=Inches(1.0),
        size=17, color=TEXT_SECONDARY, font="Source Serif 4", italic=True,
    )
    bullets = [
        ("Geometric, deterministic, no LLM per query.",
         "Spectrum positions are pure dot products.  Same query renders the same way across reloads."),
        ("Named axes, not statistical leftovers.",
         "UMAP and PCA give 'directions of max variance' that mean nothing.  Our axes are hand-curated human concepts."),
        ("Cross-surface synchronization.",
         "Hover Aaron Dixon in one chart, watch him land at different coordinates in the other three.  The shift is what the structure means."),
        ("Provenance and uncertainty are first-class metadata.",
         "Every chunk knows its audit tier and its LoC catalog URL.  Citation-grade output by construction."),
    ]
    add_bullets(
        s, bullets,
        left=Inches(0.95), top=Inches(3.3),
        width=Inches(11.5), height=Inches(3.5),
        size=15, gap=0.85,
    )
    return s


def slide_12_steering_hierarchy():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "Steering Document Hierarchy")
    add_title(s, "Six tiers.  Read top-down, build out from there.")
    rows = [
        ("Tier 1", "Orientation",
         "CLAUDE.md, README.md, PRESENTATION_REFERENCE.md"),
        ("Tier 2", "Active reference",
         "docs/*.md, rag/README.md, rag/INTERACTIVE_FEATURES_DESIGN.md, mcp-server/USAGE_GUIDE.md"),
        ("Tier 3", "Lessons learned",
         "lessons_learned.md, docs/RAG_SUBSTRATE_DECISION.md"),
        ("Tier 4", "Demo prep",
         "rag/CONFERENCE_PREP.md, rag/DEMO_SCRIPT.md"),
        ("Tier 5", "Provenance / historical record",
         "transcripts/AUDIT_TRAIL.md, CLEANED_TRANSCRIPTS_REVIEW.md, loc_healing/COVERAGE_REPORT.md"),
        ("Tier 6", "Deprecated",
         "docs/WEAVIATE_INTEGRATION_DESIGN.md  (substrate pivoted 2026-05-22)"),
    ]
    add_table(
        s,
        ["Tier", "Purpose", "Documents"],
        rows,
        left=Inches(0.95), top=Inches(2.05),
        width=Inches(11.45), height=Inches(4.2),
    )
    add_text(
        s,
        "Full hierarchy + 'when to read what' cheat sheet:  STEERING_DOCS.md at the project root.",
        left=Inches(0.95), top=Inches(6.45),
        width=Inches(11.5), height=Inches(0.4),
        size=14, color=RED, font="Inter", bold=True, italic=True,
    )
    return s


def slide_13_when_to_read():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "When to read what")
    add_title(s, "Most-common-task cheat sheet.")
    rows = [
        ("Start contributing for the first time", "CLAUDE.md"),
        ("Brief an external stakeholder", "PRESENTATION_REFERENCE.md  ->  lessons_learned.md"),
        ("Demo the live site", "rag/DEMO_SCRIPT.md"),
        ("Edit audit overlays", "CLAUDE.md (audit discipline)  +  AUDIT_TRAIL.md"),
        ("Ingest a new transcript", "transcripts/ingestion/README.md"),
        ("Touch the RAG ingest / retrieval code", "rag/README.md"),
        ("Build a new interactive surface", "rag/INTERACTIVE_FEATURES_DESIGN.md"),
        ("Deploy to staging or production", "docs/DEPLOYMENT.md  +  rag/OPERATIONS.md"),
        ("Touch the MCP server", "mcp-server/README.md  +  mcp-server/USAGE_GUIDE.md"),
        ("Touch styling / colors / accessibility", "docs/ACCESSIBILITY.md  +  CLAUDE.md (writing rules)"),
        ("Show LoC / Smithsonian reviewer the rigor", "AUDIT_TRAIL.md  +  loc_healing/COVERAGE_REPORT.md"),
    ]
    add_table(
        s,
        ["You are about to...", "Read first"],
        rows,
        left=Inches(0.95), top=Inches(2.05),
        width=Inches(11.45), height=Inches(4.7),
    )
    return s


def slide_14_next():
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    add_section_label(s, "What is deferred  /  what is next")
    add_title(s, "Known follow-on work.")
    bullets = [
        ("5 spelling-discrepancy entries.",
         "Alternative-spelling re-search against LoC catalog.  Recoverable; small, targeted."),
        ("~92,000 SME-flagged divergences.",
         "Future targeted pass with model classification.  Most editorial; a minority real corrections we declined to auto-apply."),
        ("Forced-alignment timing improvement.",
         "WhisperX or Montreal Forced Aligner against LoC audio.  Tightens per-clip precision below the 5-second cue grid."),
        ("MCP server deploy to Fly.io.",
         "Code verified locally.  Awaiting flyctl auth login + fly deploy."),
        ("Pipeline run on the 135 raw transcripts.",
         "Dual-scoring + citation audit.  ~$5.40 estimated cost across the full corpus."),
        ("Firebase Cloud Functions deploy.",
         "Requires Blaze billing on the new project + the OPENAI_API_KEY secret."),
    ]
    add_bullets(
        s, bullets,
        left=Inches(0.95), top=Inches(2.05),
        width=Inches(11.5), height=Inches(4.5),
        size=14, gap=0.75,
    )
    return s


def slide_15_close():
    s = prs.slides.add_slide(BLANK)
    add_background(s, color=BG_TITLE)

    # Decorative right bar
    add_accent_bar(
        s,
        left=prs.slide_width - Inches(0.12), top=Inches(0),
        height=prs.slide_height, color=RED,
    )

    add_text(
        s, "WHAT TO TAKE AWAY",
        left=Inches(0.95), top=Inches(0.9),
        width=Inches(11), height=Inches(0.4),
        size=13, bold=True, color=RED, font="Chivo Mono",
    )

    add_text(
        s, "We built an audit instrument,",
        left=Inches(0.95), top=Inches(1.55),
        width=Inches(12), height=Inches(1.1),
        size=44, bold=True, color=TEXT_PRIMARY, font="Inter",
    )
    add_text(
        s, "not a clean-up script.",
        left=Inches(0.95), top=Inches(2.4),
        width=Inches(12), height=Inches(1.1),
        size=44, bold=True, color=RED, font="Inter",
    )

    add_text(
        s,
        "Eight internal passes plus a Library of Congress healing pass against the "
        "canonical archive.  Conservative-first-pass discipline so we do not auto-introduce "
        "new errors.  Every passage carries its audit tier and its LoC catalog deep-link.",
        left=Inches(0.95), top=Inches(3.75),
        width=Inches(11.5), height=Inches(2.0),
        size=17, color=TEXT_SECONDARY, font="Source Serif 4",
    )

    add_text(
        s,
        "civil-rights-staging.netlify.app  /  /rag-explore",
        left=Inches(0.95), top=Inches(5.4),
        width=Inches(11.5), height=Inches(0.5),
        size=14, color=TEXT_MUTED, font="Chivo Mono",
    )

    # Big Questions.
    line = s.shapes.add_connector(
        1, Inches(0.95), Inches(6.1), Inches(4), Inches(6.1)
    )
    line.line.color.rgb = RED
    line.line.width = Pt(2)
    add_text(
        s, "Questions.",
        left=Inches(0.95), top=Inches(6.2),
        width=Inches(11.5), height=Inches(0.9),
        size=44, bold=True, color=RED, font="Inter",
    )
    return s


renderers = [
    slide_01_title,
    slide_02_problem,
    slide_03_error_examples,
    slide_04_audit_cascade,
    slide_05_loc_pass,
    slide_06_coverage,
    slide_07_conservative,
    slide_08_challenges,
    slide_09_what_it_enables,
    slide_10_demos,
    slide_mcp,
    slide_11_philosophy,
    slide_12_steering_hierarchy,
    slide_13_when_to_read,
    slide_14_next,
    slide_15_close,
]


total = len(renderers)
for i, r in enumerate(renderers, start=1):
    slide = r()
    if i != 1 and i != total:
        add_footer(slide, i, total)


out_path = Path("wwu_presentation_2026-05-27_dark.pptx")
try:
    prs.save(out_path)
except PermissionError:
    # Caller has the file open in PowerPoint; bump the filename so the
    # write still succeeds and we don't clobber their in-progress edits.
    import datetime
    stamp = datetime.datetime.now().strftime("%H%M%S")
    out_path = Path(f"wwu_presentation_2026-05-27_dark_{stamp}.pptx")
    prs.save(out_path)
print(f"Wrote {out_path.resolve()}")
print(f"  Total slides: {total}")
